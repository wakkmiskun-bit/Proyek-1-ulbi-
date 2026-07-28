import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx library not found. Please install it with 'pip install python-pptx'.")
    sys.exit(1)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Theme Colors (Original Navy/Pink Template)
    COLOR_PRIMARY = RGBColor(233, 30, 99)     # Deep Pink (#e91e63)
    COLOR_SECONDARY = RGBColor(240, 98, 146) # Light Pink (#f06292)
    COLOR_DARK = RGBColor(26, 26, 46)        # Dark Navy (#1a1a2e)
    COLOR_CARD_BG = RGBColor(245, 247, 250)  # Card Background (#f5f7fa)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT = RGBColor(40, 40, 40)
    COLOR_MUTED = RGBColor(120, 120, 120)
    COLOR_BORDER = RGBColor(226, 232, 240)

    blank_layout = prs.slide_layouts[6]

    # Image paths
    img_dir = os.path.join(os.path.dirname(__file__), "images")
    img_landing = os.path.join(img_dir, "landing_page.jpg")
    img_kanban = os.path.join(img_dir, "student_kanban.jpg")
    img_admin = os.path.join(img_dir, "admin_dashboard.jpg")
    img_wa = os.path.join(img_dir, "wa_reminder.jpg")

    # Helper function to add a card
    def add_card(slide, left, top, width, height, bg_color=COLOR_CARD_BG, border_color=COLOR_BORDER):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()
        return card

    # Original Header/Footer Template from generate_pptx.py
    def add_slide_header_and_footer(slide, title_text, category_text="TASKMATE - PRESENTASI SIDANG PROYEK 1"):
        # Top Dark Navy Header Band
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_DARK
        shape.line.fill.background()

        # Pink Accent Strip
        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.08))
        strip.fill.solid()
        strip.fill.fore_color.rgb = COLOR_PRIMARY
        strip.line.fill.background()

        # Category Text
        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        tf_cat.margin_left = tf_cat.margin_right = tf_cat.margin_top = tf_cat.margin_bottom = 0
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_SECONDARY
        p_cat.font.name = "Arial"

        # Title Text
        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11.7), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        tf_title.margin_left = tf_title.margin_right = tf_title.margin_top = tf_title.margin_bottom = 0
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE
        p_title.font.name = "Arial"

        # Footer Text
        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.35))
        tf_foot = footer_box.text_frame
        tf_foot.margin_left = tf_foot.margin_right = tf_foot.margin_top = tf_foot.margin_bottom = 0
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "TaskMate v2.0 | D4 Teknik Informatika ULBI 2026 | Kelompok 22: M. Ilham Habiballah & Gianjar Nugraha"
        p_foot.font.size = Pt(9.5)
        p_foot.font.color.rgb = COLOR_MUTED
        p_foot.font.name = "Arial"

        # Footer Divider Line
        foot_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.02))
        foot_line.fill.solid()
        foot_line.fill.fore_color.rgb = COLOR_SECONDARY
        foot_line.line.fill.background()

    # Laptop Mockup helper
    def add_laptop_mockup(slide, left, top, width, height, image_path):
        bezel_w = width
        bezel_h = height * 0.88
        bezel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, bezel_w, bezel_h)
        bezel.fill.solid()
        bezel.fill.fore_color.rgb = COLOR_DARK
        bezel.line.fill.background()
        
        screen_w = bezel_w * 0.94
        screen_h = bezel_h * 0.88
        screen_l = left + (bezel_w - screen_w) / 2
        screen_t = top + (bezel_h - screen_h) / 2
        
        if os.path.exists(image_path):
            pic = slide.shapes.add_picture(image_path, screen_l, screen_t, width=screen_w, height=screen_h)
            pic.line.color.rgb = COLOR_PRIMARY
            pic.line.width = Pt(1)
            
        kb_w = bezel_w * 1.08
        kb_h = height * 0.12
        kb_l = left - (kb_w - bezel_w) / 2
        kb_t = top + bezel_h - Inches(0.03)
        
        kb = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, kb_l, kb_t, kb_w, kb_h)
        kb.fill.solid()
        kb.fill.fore_color.rgb = RGBColor(200, 200, 200)
        kb.line.fill.background()
        
        tp_w = kb_w * 0.18
        tp_h = kb_h * 0.6
        tp_l = kb_l + (kb_w - tp_w) / 2
        tp_t = kb_t + (kb_h - tp_h) / 2
        
        tp = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, tp_l, tp_t, tp_w, tp_h)
        tp.fill.solid()
        tp.fill.fore_color.rgb = RGBColor(160, 160, 160)
        tp.line.fill.background()

    # Phone Mockup helper
    def add_phone_mockup(slide, left, top, width, height, image_path):
        bezel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        bezel.fill.solid()
        bezel.fill.fore_color.rgb = COLOR_DARK
        bezel.line.fill.background()
        
        screen_w = width * 0.92
        screen_h = height * 0.94
        screen_l = left + (width - screen_w) / 2
        screen_t = top + (height - screen_h) / 2
        
        if os.path.exists(image_path):
            pic = slide.shapes.add_picture(image_path, screen_l, screen_t, width=screen_w, height=screen_h)
            pic.line.color.rgb = COLOR_PRIMARY
            pic.line.width = Pt(1)
            
        notch_w = width * 0.35
        notch_h = height * 0.03
        notch_l = left + (width - notch_w) / 2
        notch_t = top + Inches(0.05)
        
        notch = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, notch_l, notch_t, notch_w, notch_h)
        notch.fill.solid()
        notch.fill.fore_color.rgb = RGBColor(0, 0, 0)
        notch.line.fill.background()

    # Helper function to add circles with characters/emojis
    def add_circle_badge(slide, left, top, size, bg_color, text, font_size=16):
        badge = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
        badge.fill.solid()
        badge.fill.fore_color.rgb = bg_color
        badge.line.fill.background()
        
        tf = badge.text_frame
        tf.word_wrap = False
        tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.bold = True
        p.font.color.rgb = COLOR_WHITE
        p.font.name = "Segoe UI"
        p.alignment = PP_ALIGN.CENTER
        return badge

    # =========================================================================
    # SLIDE 1: COVER (Navy Theme matching original presentation)
    # =========================================================================
    slide1 = prs.slides.add_slide(blank_layout)
    
    # Solid Navy background
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_DARK
    bg1.line.fill.background()

    # Pink decorative strip on the left edge
    dec = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    dec.fill.solid()
    dec.fill.fore_color.rgb = COLOR_PRIMARY
    dec.line.fill.background()

    # Title & Subtitle block
    tbox = slide1.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(6.8), Inches(2.5))
    tf = tbox.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf_url = 0

    p = tf.paragraphs[0]
    p.text = "TASKMATE"
    p.font.size = Pt(54)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = "Sistem Manajemen Tugas & Produktivitas Berbasis Web"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.font.name = "Arial"
    p2.space_before = Pt(6)

    p3 = tf.add_paragraph()
    p3.text = "Laporan Akhir Presentasi Sidang Proyek 1"
    p3.font.size = Pt(15)
    p3.font.color.rgb = COLOR_SECONDARY
    p3.font.name = "Arial"
    p3.space_before = Pt(8)

    # Embed Landing Page UI Image on Cover (Right Column)
    if os.path.exists(img_landing):
        slide1.shapes.add_picture(img_landing, Inches(7.8), Inches(0.8), width=Inches(4.8), height=Inches(3.0))
        # Add pink border
        border1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(7.8), Inches(0.8), Inches(4.8), Inches(3.0))
        border1.fill.background()
        border1.line.color.rgb = COLOR_PRIMARY
        border1.line.width = Pt(1.5)

    # Cards for Presenters (Kelompok 22)
    add_card(slide1, Inches(0.8), Inches(4.0), Inches(5.6), Inches(2.7), bg_color=RGBColor(35, 35, 60), border_color=COLOR_PRIMARY)
    tb_pres = slide1.shapes.add_textbox(Inches(1.0), Inches(4.15), Inches(5.2), Inches(2.4))
    tf_p = tb_pres.text_frame
    tf_p.word_wrap = True
    tf_p.margin_left = tf_p.margin_right = tf_p.margin_top = 0

    p_p1 = tf_p.paragraphs[0]
    p_p1.text = "DISUSUN OLEH (KELOMPOK 22):"
    p_p1.font.size = Pt(12)
    p_p1.font.bold = True
    p_p1.font.color.rgb = COLOR_SECONDARY

    p_p2 = tf_p.add_paragraph()
    p_p2.text = "1. Muhammad Ilham Habiballah (NPM: 714250003)"
    p_p2.font.size = Pt(13)
    p_p2.font.bold = True
    p_p2.font.color.rgb = COLOR_WHITE
    p_p2.space_before = Pt(8)

    p_p3 = tf_p.add_paragraph()
    p_p3.text = "2. Gianjar Nugraha (NPM: 714250007)"
    p_p3.font.size = Pt(13)
    p_p3.font.bold = True
    p_p3.font.color.rgb = COLOR_WHITE
    p_p3.space_before = Pt(6)

    p_p4 = tf_p.add_paragraph()
    p_p4.text = "Program Studi D4 Teknik Informatika\nUniversitas Logistik & Bisnis Internasional"
    p_p4.font.size = Pt(11)
    p_p4.font.color.rgb = COLOR_MUTED
    p_p4.space_before = Pt(12)

    # Card for Supervisor
    add_card(slide1, Inches(6.8), Inches(4.0), Inches(5.7), Inches(2.7), bg_color=RGBColor(35, 35, 60), border_color=COLOR_SECONDARY)
    tb_dos = slide1.shapes.add_textbox(Inches(7.0), Inches(4.15), Inches(5.3), Inches(2.4))
    tf_d = tb_dos.text_frame
    tf_d.word_wrap = True
    tf_d.margin_left = tf_d.margin_right = tf_d.margin_top = 0

    p_d1 = tf_d.paragraphs[0]
    p_d1.text = "DOSEN PEMBIMBING & INSTITUSI:"
    p_d1.font.size = Pt(12)
    p_d1.font.bold = True
    p_d1.font.color.rgb = COLOR_SECONDARY

    p_d2 = tf_d.add_paragraph()
    p_d2.text = "Pembimbing: Cahyo Prianto, S.Pd., M.T., CDSP, SFPC"
    p_d2.font.size = Pt(13)
    p_d2.font.bold = True
    p_d2.font.color.rgb = COLOR_WHITE
    p_d2.space_before = Pt(8)

    p_d3 = tf_d.add_paragraph()
    p_d3.text = "Tahun Akademik: 2025/2026"
    p_d3.font.size = Pt(12.5)
    p_d3.font.color.rgb = COLOR_WHITE
    p_d3.space_before = Pt(6)

    p_d4 = tf_d.add_paragraph()
    p_d4.text = "D4 Teknik Informatika ULBI 2026  |  TASKMATE.ULBI.AC.ID"
    p_d4.font.size = Pt(11)
    p_d4.font.color.rgb = COLOR_MUTED
    p_d4.space_before = Pt(20)

    # =========================================================================
    # SLIDE 2: LATAR BELAKANG & TUJUAN
    # =========================================================================
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(slide2, "Latar Belakang & Tujuan Solutif")

    # Left Card: Masalah
    add_card(slide2, Inches(0.8), Inches(1.3), Inches(5.7), Inches(5.2), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
    add_circle_badge(slide2, Inches(1.1), Inches(1.5), Inches(0.45), COLOR_PRIMARY, "!", font_size=14)
    
    masalah_title = slide2.shapes.add_textbox(Inches(1.65), Inches(1.5), Inches(4.5), Inches(0.45))
    tf_mt = masalah_title.text_frame
    tf_mt.margin_left = tf_mt.margin_right = tf_mt.margin_top = 0
    p_mt = tf_mt.paragraphs[0]
    p_mt.text = "Identifikasi Masalah"
    p_mt.font.size = Pt(16)
    p_mt.font.bold = True
    p_mt.font.color.rgb = COLOR_DARK

    points_masalah = [
        "Beban aktivitas akademik mahasiswa D4 TI ULBI yang sangat padat.",
        "Pencatatan tugas harian/praktikum konvensional bersifat pasif.",
        "Ketiadaan pengingat proaktif yang terintegrasi (forgetfulness).",
        "Kesulitan menyusun skala prioritas pengerjaan tugas secara dinamis.",
        "Pengelola (Admin) kesulitan memantau produktivitas tugas mahasiswa."
    ]
    
    m_list_box = slide2.shapes.add_textbox(Inches(1.1), Inches(2.1), Inches(5.1), Inches(2.3))
    tf_ml = m_list_box.text_frame
    tf_ml.word_wrap = True
    tf_ml.margin_left = tf_ml.margin_right = tf_ml.margin_top = 0
    for idx, pt in enumerate(points_masalah):
        p = tf_ml.add_paragraph() if idx > 0 else tf_ml.paragraphs[0]
        p.text = f"{idx+1}   {pt}"
        p.font.size = Pt(10.5)
        p.font.color.rgb = COLOR_TEXT
        p.space_after = Pt(5)
        
        run = p.runs[0]
        run.font.bold = True
        run.font.color.rgb = COLOR_PRIMARY

    # Small application preview inside Left Card using laptop mockup
    add_laptop_mockup(slide2, Inches(1.2), Inches(4.5), Inches(2.3), Inches(1.5), img_landing)

    # Text next to preview
    desc_prev = slide2.shapes.add_textbox(Inches(3.7), Inches(4.8), Inches(2.5), Inches(1.2))
    tf_dp = desc_prev.text_frame
    tf_dp.word_wrap = True
    tf_dp.margin_left = tf_dp.margin_right = tf_dp.margin_top = 0
    p_dp = tf_dp.paragraphs[0]
    p_dp.text = "TaskMate UI / UX\nModern Dashboard dengan visual manajemen responsif."
    p_dp.font.size = Pt(10)
    p_dp.font.bold = True
    p_dp.font.color.rgb = COLOR_MUTED

    # Right Card: Tujuan
    add_card(slide2, Inches(6.8), Inches(1.3), Inches(5.7), Inches(5.2), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
    add_circle_badge(slide2, Inches(7.1), Inches(1.5), Inches(0.45), COLOR_DARK, "🎯", font_size=12)

    tujuan_title = slide2.shapes.add_textbox(Inches(7.65), Inches(1.5), Inches(4.5), Inches(0.45))
    tf_tt = tujuan_title.text_frame
    tf_tt.margin_left = tf_tt.margin_right = tf_tt.margin_top = 0
    p_tt = tf_tt.paragraphs[0]
    p_tt.text = "Tujuan Solutif"
    p_tt.font.size = Pt(16)
    p_tt.font.bold = True
    p_tt.font.color.rgb = COLOR_DARK

    tujuan_desc = slide2.shapes.add_textbox(Inches(7.1), Inches(2.1), Inches(5.1), Inches(0.9))
    tf_td = tujuan_desc.text_frame
    tf_td.word_wrap = True
    tf_td.margin_left = tf_td.margin_right = tf_td.margin_top = 0
    p_td = tf_td.paragraphs[0]
    p_td.text = "Membangun sistem informasi TaskMate berbasis web untuk mengelola dan memonitor tugas akademik mahasiswa secara efisien, proaktif, dan real-time."
    p_td.font.size = Pt(10.5)
    p_td.font.color.rgb = COLOR_TEXT
    p_td.space_after = Pt(10)

    # 4 Grid Feature boxes inside Right Card
    f_boxes = [
        ("Visual Kanban Board", "4 Papan Kolom Tugas", COLOR_PRIMARY),
        ("WhatsApp Gateway", "Fonnte Auto Reminder", COLOR_PRIMARY),
        ("Multi-Role Access", "Admin & Mahasiswa Guard", COLOR_PRIMARY),
        ("Live Monitoring", "Statistik & Audit Trail", COLOR_PRIMARY)
    ]
    
    positions = [
        (Inches(7.1), Inches(3.1)), # Top Left
        (Inches(9.75), Inches(3.1)), # Top Right
        (Inches(7.1), Inches(4.25)), # Bottom Left
        (Inches(9.75), Inches(4.25))  # Bottom Right
    ]

    for i, (title, desc, color) in enumerate(f_boxes):
        pos_x, pos_y = positions[i]
        box = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, pos_x, pos_y, Inches(2.45), Inches(0.95))
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        
        tf_b = box.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = tf_b.margin_bottom = Inches(0.05)
        
        p1 = tf_b.paragraphs[0]
        p1.text = title
        p1.font.size = Pt(10)
        p1.font.bold = True
        p1.font.color.rgb = COLOR_WHITE
        p1.alignment = PP_ALIGN.CENTER
        
        p2 = tf_b.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(8.5)
        p2.font.color.rgb = COLOR_SECONDARY
        p2.alignment = PP_ALIGN.CENTER
        p2.space_before = Pt(4)

    # Output badge at bottom right
    clip_shape = slide2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.1), Inches(5.4), Inches(5.1), Inches(0.85))
    clip_shape.fill.solid()
    clip_shape.fill.fore_color.rgb = COLOR_CARD_BG
    clip_shape.line.color.rgb = COLOR_BORDER
    clip_shape.line.width = Pt(1.0)

    tf_cs = clip_shape.text_frame
    tf_cs.word_wrap = True
    p_cs = tf_cs.paragraphs[0]
    p_cs.text = "🎯 Output Akhir: Produktivitas belajar mahasiswa meningkat, angka keterlambatan tugas dapat diturunkan secara drastis."
    p_cs.font.size = Pt(9.5)
    p_cs.font.bold = True
    p_cs.font.color.rgb = COLOR_PRIMARY
    p_cs.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 3: FITUR UNGGULAN
    # =========================================================================
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(slide3, "Fitur Unggulan Aplikasi")

    # 4 Cards in 2x2 grid
    card_data = [
        ("Multi-Role Authorization", 
         "Pemisahan hak akses aman menggunakan Guard Laravel (Mahasiswa untuk workspace tugas terisolasi, Admin untuk dasbor pengawasan & monitoring global).",
         "🔐", Inches(0.8), Inches(1.3)),
         
        ("Visual Kanban Board", 
         "Manajemen tugas dengan 4 kolom visual (To Do, Doing, Review, Done) didukung SortableJS drag-and-drop dan auto-done 100% checklist sub-tugas.",
         "📋", Inches(6.8), Inches(1.3)),
         
        ("WhatsApp Gateway Reminder", 
         "Pengingat otomatis tenggat waktu H-5 & H-2 sebelum batas waktu berakhir secara terjadwal (Laravel Cron Scheduler) langsung ke WhatsApp via Fonnte API.",
         "💬", Inches(0.8), Inches(3.9)),
         
        ("Live Monitoring & Audit Log", 
         "Dasbor analitik monitoring progres board mahasiswa secara real-time dan tracking audit trail log aktivitas sistem untuk transparansi.",
         "📊", Inches(6.8), Inches(3.9))
    ]

    for title, desc, icon, x, y in card_data:
        add_card(slide3, x, y, Inches(5.7), Inches(2.3), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
        
        # Circular Badge on left of card
        add_circle_badge(slide3, x + Inches(0.3), y + Inches(0.45), Inches(1.0), COLOR_PRIMARY, icon, font_size=20)
        
        # Text Frame on right of card
        txt_box = slide3.shapes.add_textbox(x + Inches(1.5), y + Inches(0.25), Inches(3.9), Inches(1.8))
        tf_tx = txt_box.text_frame
        tf_tx.word_wrap = True
        tf_tx.margin_left = tf_tx.margin_right = tf_tx.margin_top = tf_tx.margin_bottom = 0
        
        p_title = tf_tx.paragraphs[0]
        p_title.text = title
        p_title.font.size = Pt(14)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_PRIMARY
        p_title.space_after = Pt(6)
        
        p_desc = tf_tx.add_paragraph()
        p_desc.text = desc
        p_desc.font.size = Pt(10)
        p_desc.font.color.rgb = COLOR_TEXT
        p_desc.line_spacing = 1.15

    # =========================================================================
    # SLIDE 4: ALUR KERJA APLIKASI
    # =========================================================================
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(slide4, "Alur Kerja Operasional")

    # Column Left: Mahasiswa
    add_card(slide4, Inches(0.8), Inches(1.3), Inches(5.7), Inches(4.5), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
    
    col1_title = slide4.shapes.add_textbox(Inches(1.0), Inches(1.45), Inches(5.3), Inches(0.45))
    tf_c1t = col1_title.text_frame
    tf_c1t.margin_left = tf_c1t.margin_right = tf_c1t.margin_top = 0
    p_c1t = tf_c1t.paragraphs[0]
    p_c1t.text = "👤 PENGELOLAAN TUGAS (MAHASISWA)"
    p_c1t.font.size = Pt(13)
    p_c1t.font.bold = True
    p_c1t.font.color.rgb = COLOR_PRIMARY
    p_c1t.alignment = PP_ALIGN.CENTER

    steps_mhs = [
        ("Pilih Kriteria / Modul", "Mahasiswa input tugas akademik beserta prioritas & sub-tugas."),
        ("Kanban Lifecycle", "Menyeret kartu tugas (To Do -> Doing -> Review -> Done) via AJAX."),
        ("Auto-Done Checklist", "Tugas otomatis bergeser ke 'Done' ketika sub-tugas mencapai 100%."),
        ("Terima WhatsApp Reminder", "Notifikasi WA Gateway diterima otomatis H-5 & H-2 sebelum deadline.")
    ]

    for idx, (title, desc) in enumerate(steps_mhs):
        y_pos = Inches(2.0) + idx * Inches(0.85)
        # Circle badge for number
        add_circle_badge(slide4, Inches(1.1), y_pos, Inches(0.45), COLOR_PRIMARY, str(idx+1), font_size=12)
        
        # Text block
        tb = slide4.shapes.add_textbox(Inches(1.7), y_pos - Inches(0.05), Inches(4.6), Inches(0.8))
        tf_b = tb.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = 0
        
        p1 = tf_b.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(11)
        p1.font.color.rgb = COLOR_DARK
        
        p2 = tf_b.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = COLOR_TEXT
        p2.space_before = Pt(2)

    # Column Right: Admin
    add_card(slide4, Inches(6.8), Inches(1.3), Inches(5.7), Inches(4.5), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
    
    col2_title = slide4.shapes.add_textbox(Inches(7.0), Inches(1.45), Inches(5.3), Inches(0.45))
    tf_c2t = col2_title.text_frame
    tf_c2t.margin_left = tf_c2t.margin_right = tf_c2t.margin_top = 0
    p_c2t = tf_c2t.paragraphs[0]
    p_c2t.text = "🔐 MONITORING & VERIFIKASI (ADMIN)"
    p_c2t.font.size = Pt(13)
    p_c2t.font.bold = True
    p_c2t.font.color.rgb = COLOR_DARK
    p_c2t.alignment = PP_ALIGN.CENTER

    steps_admin = [
        ("Akses Dashboard Admin", "Memantau performa, grafik tugas, dan statistik keaktifan global."),
        ("Live Board Monitoring", "Mengecek progres Kanban board live masing-masing mahasiswa."),
        ("Audit Log Activity", "Memverifikasi log audit trail digital transaksi tindakan mahasiswa."),
        ("Trigger WhatsApp Scheduler", "Memicu pengiriman reminder terjadwal lewat CLI command Fonnte.")
    ]

    for idx, (title, desc) in enumerate(steps_admin):
        y_pos = Inches(2.0) + idx * Inches(0.85)
        # Circle badge for number
        add_circle_badge(slide4, Inches(7.1), y_pos, Inches(0.45), COLOR_DARK, str(idx+1), font_size=12)
        
        # Text block
        tb = slide4.shapes.add_textbox(Inches(7.7), y_pos - Inches(0.05), Inches(4.6), Inches(0.8))
        tf_b = tb.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = 0
        
        p1 = tf_b.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(11)
        p1.font.color.rgb = COLOR_DARK
        
        p2 = tf_b.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(9.5)
        p2.font.color.rgb = COLOR_TEXT
        p2.space_before = Pt(2)

    # Bottom capsule (outcome)
    outcome_cap = slide4.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(2.5), Inches(6.0), Inches(8.333), Inches(0.45))
    outcome_cap.fill.solid()
    outcome_cap.fill.fore_color.rgb = COLOR_DARK
    outcome_cap.line.fill.background()

    tf_oc = outcome_cap.text_frame
    tf_oc.word_wrap = True
    p_oc = tf_oc.paragraphs[0]
    p_oc.text = "TUGAS AKADEMIK SELESAI & TERMONITOR SECARA EFEKTIF"
    p_oc.font.size = Pt(11.5)
    p_oc.font.bold = True
    p_oc.font.color.rgb = COLOR_WHITE
    p_oc.font.name = "Arial"
    p_oc.alignment = PP_ALIGN.CENTER

    # =========================================================================
    # SLIDE 5: KESIMPULAN & WHATSAPP GATEWAY
    # =========================================================================
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(slide5, "Kesimpulan & Integrasi WhatsApp")

    # Left Column: Kesimpulan
    add_card(slide5, Inches(0.8), Inches(1.3), Inches(6.0), Inches(5.2), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
    add_circle_badge(slide5, Inches(1.1), Inches(1.5), Inches(0.45), COLOR_PRIMARY, "📋", font_size=12)

    kes_title = slide5.shapes.add_textbox(Inches(1.65), Inches(1.5), Inches(5.0), Inches(0.45))
    tf_kt = kes_title.text_frame
    tf_kt.margin_left = tf_kt.margin_right = tf_kt.margin_top = 0
    p_kt = tf_kt.paragraphs[0]
    p_kt.text = "Kesimpulan Implementasi"
    p_kt.font.size = Pt(16)
    p_kt.font.bold = True
    p_kt.font.color.rgb = COLOR_DARK

    conclusions = [
        ("Mendigitalisasi Papan Kerja", "Mengubah sistem pencatatan manual menjadi visual Kanban Board terpusat."),
        ("Efektivitas WhatsApp Gateway", "Notifikasi proaktif H-5 & H-2 berhasil meminimalkan keterlambatan pengumpulan."),
        ("Kemudahan Pengawasan", "Dasbor admin mempermudah monitor progres tugas real-time & audit activity trail."),
        ("Pusat Bantuan Terpadu", "Fitur Direct Help Center WA mempercepat penyelesaian aduan kendala dari mahasiswa.")
    ]

    for idx, (title, desc) in enumerate(conclusions):
        y_pos = Inches(2.1) + idx * Inches(1.1)
        # Circle badge for number
        add_circle_badge(slide5, Inches(1.1), y_pos, Inches(0.42), COLOR_PRIMARY, str(idx+1), font_size=11)
        
        # Text block
        tb = slide5.shapes.add_textbox(Inches(1.65), y_pos - Inches(0.05), Inches(5.0), Inches(1.0))
        tf_b = tb.text_frame
        tf_b.word_wrap = True
        tf_b.margin_left = tf_b.margin_right = tf_b.margin_top = 0
        
        p1 = tf_b.paragraphs[0]
        p1.text = title
        p1.font.bold = True
        p1.font.size = Pt(11.5)
        p1.font.color.rgb = COLOR_DARK
        
        p2 = tf_b.add_paragraph()
        p2.text = desc
        p2.font.size = Pt(10)
        p2.font.color.rgb = COLOR_TEXT
        p2.space_before = Pt(2)

    # Right Column: WhatsApp Gateway Mockup
    add_card(slide5, Inches(7.0), Inches(1.3), Inches(5.5), Inches(5.2), bg_color=COLOR_WHITE, border_color=COLOR_BORDER)
    add_circle_badge(slide5, Inches(7.3), Inches(1.5), Inches(0.45), COLOR_PRIMARY, "💬", font_size=12)

    wa_title = slide5.shapes.add_textbox(Inches(7.85), Inches(1.5), Inches(4.5), Inches(0.45))
    tf_wt = wa_title.text_frame
    tf_wt.margin_left = tf_wt.margin_right = tf_wt.margin_top = 0
    p_wt = tf_wt.paragraphs[0]
    p_wt.text = "Sistem Notifikasi WhatsApp"
    p_wt.font.size = Pt(15)
    p_wt.font.bold = True
    p_wt.font.color.rgb = COLOR_DARK

    # WhatsApp image screenshot using phone mockup
    add_phone_mockup(slide5, Inches(8.3), Inches(2.0), Inches(2.9), Inches(3.8), img_wa)

    # Gold/Orange badge at the bottom of Right Card (replicates HKI badge on SILADATA Slide 5)
    badge_capsule = slide5.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(7.3), Inches(5.95), Inches(4.9), Inches(0.45))
    badge_capsule.fill.solid()
    badge_capsule.fill.fore_color.rgb = COLOR_PRIMARY
    badge_capsule.line.fill.background()

    tf_bc = badge_capsule.text_frame
    tf_bc.word_wrap = True
    p_bc = tf_bc.paragraphs[0]
    p_bc.text = "WHATSAPP GATEWAY INTEGRATED - FONNTE API VERIFIED"
    p_bc.font.size = Pt(9.5)
    p_bc.font.bold = True
    p_bc.font.color.rgb = COLOR_WHITE
    p_bc.font.name = "Arial"
    p_bc.alignment = PP_ALIGN.CENTER

    # Save the output pptx
    output_path = "TaskMate_Presentasi_5_Slide.pptx"
    prs.save(output_path)
    print(f"Presentation saved successfully to: {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
