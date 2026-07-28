import sys
import os

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
except ImportError:
    print("python-pptx library not found. Please install it with 'pip install python-pptx'.")
    sys.exit(1)

def create_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    # Theme Colors
    COLOR_PRIMARY = RGBColor(233, 30, 99)     # Deep Pink (#e91e63)
    COLOR_SECONDARY = RGBColor(240, 98, 146) # Light Pink (#f06292)
    COLOR_DARK = RGBColor(26, 26, 46)        # Dark Navy (#1a1a2e)
    COLOR_CARD = RGBColor(245, 247, 250)     # Card Background (#f5f7fa)
    COLOR_WHITE = RGBColor(255, 255, 255)
    COLOR_TEXT = RGBColor(40, 40, 40)
    COLOR_MUTED = RGBColor(120, 120, 120)

    blank_layout = prs.slide_layouts[6]

    # Image paths
    img_dir = os.path.join(os.path.dirname(__file__), "images")
    img_landing = os.path.join(img_dir, "landing_page.jpg")
    img_kanban = os.path.join(img_dir, "student_kanban.jpg")
    img_admin = os.path.join(img_dir, "admin_dashboard.jpg")

    def add_slide_header_and_footer(slide, title_text, category_text="TASKMATE - PRESENTASI SIDANG PROYEK 1"):
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.1))
        shape.fill.solid()
        shape.fill.fore_color.rgb = COLOR_DARK
        shape.line.fill.background()

        strip = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.1), Inches(13.333), Inches(0.08))
        strip.fill.solid()
        strip.fill.fore_color.rgb = COLOR_PRIMARY
        strip.line.fill.background()

        cat_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.12), Inches(11.7), Inches(0.3))
        tf_cat = cat_box.text_frame
        tf_cat.word_wrap = True
        p_cat = tf_cat.paragraphs[0]
        p_cat.text = category_text.upper()
        p_cat.font.size = Pt(10)
        p_cat.font.bold = True
        p_cat.font.color.rgb = COLOR_SECONDARY
        p_cat.font.name = "Arial"

        title_box = slide.shapes.add_textbox(Inches(0.8), Inches(0.38), Inches(11.7), Inches(0.6))
        tf_title = title_box.text_frame
        tf_title.word_wrap = True
        p_title = tf_title.paragraphs[0]
        p_title.text = title_text
        p_title.font.size = Pt(22)
        p_title.font.bold = True
        p_title.font.color.rgb = COLOR_WHITE
        p_title.font.name = "Arial"

        footer_box = slide.shapes.add_textbox(Inches(0.8), Inches(7.05), Inches(11.7), Inches(0.35))
        tf_foot = footer_box.text_frame
        p_foot = tf_foot.paragraphs[0]
        p_foot.text = "TaskMate v2.0 | D4 Teknik Informatika ULBI 2026 | Kelompok 26: M. Ilham Habiballah & Gianjar Nugraha"
        p_foot.font.size = Pt(9.5)
        p_foot.font.color.rgb = COLOR_MUTED
        p_foot.font.name = "Arial"

        foot_line = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(7.0), Inches(11.7), Inches(0.02))
        foot_line.fill.solid()
        foot_line.fill.fore_color.rgb = COLOR_SECONDARY
        foot_line.line.fill.background()

    def add_card(slide, left, top, width, height, title="", bg_color=COLOR_CARD, border_color=None):
        card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        card.fill.solid()
        card.fill.fore_color.rgb = bg_color
        if border_color:
            card.line.color.rgb = border_color
            card.line.width = Pt(1.5)
        else:
            card.line.fill.background()

        if title:
            tb = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.15), width - Inches(0.4), Inches(0.5))
            tf = tb.text_frame
            tf.word_wrap = True
            p = tf.paragraphs[0]
            p.text = title
            p.font.size = Pt(13.5)
            p.font.bold = True
            p.font.color.rgb = COLOR_PRIMARY
            p.font.name = "Arial"

        return card

    # ==================== SLIDE 1: COVER ====================
    slide1 = prs.slides.add_slide(blank_layout)
    bg1 = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = COLOR_DARK
    bg1.line.fill.background()

    dec = slide1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(0.4), Inches(7.5))
    dec.fill.solid()
    dec.fill.fore_color.rgb = COLOR_PRIMARY
    dec.line.fill.background()

    tbox = slide1.shapes.add_textbox(Inches(0.8), Inches(0.8), Inches(7.2), Inches(2.2))
    tf = tbox.text_frame
    tf.word_wrap = True

    p = tf.paragraphs[0]
    p.text = "TASKMATE"
    p.font.size = Pt(46)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.font.name = "Arial"

    p2 = tf.add_paragraph()
    p2.text = "Sistem Manajemen Tugas & Produktivitas Berbasis Web"
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = COLOR_WHITE
    p2.font.name = "Arial"

    p3 = tf.add_paragraph()
    p3.text = "Laporan Akhir Presentasi Sidang Proyek 1"
    p3.font.size = Pt(15)
    p3.font.color.rgb = COLOR_SECONDARY
    p3.font.name = "Arial"
    p3.space_before = Pt(8)

    # Embed Landing Page UI Image on Cover
    if os.path.exists(img_landing):
        slide1.shapes.add_picture(img_landing, Inches(7.8), Inches(0.8), width=Inches(4.8))

    add_card(slide1, Inches(0.8), Inches(4.0), Inches(5.6), Inches(2.7), bg_color=RGBColor(35, 35, 60), border_color=COLOR_PRIMARY)
    tb_pres = slide1.shapes.add_textbox(Inches(1.0), Inches(4.1), Inches(5.2), Inches(2.5))
    tf_p = tb_pres.text_frame
    tf_p.word_wrap = True

    p = tf_p.paragraphs[0]
    p.text = "DISUSUN OLEH (KELOMPOK 26):"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY

    p = tf_p.add_paragraph()
    p.text = "1. Muhammad Ilham Habiballah (NPM: 714250003)"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_WHITE
    p.space_before = Pt(6)

    p = tf_p.add_paragraph()
    p.text = "2. Gianjar Nugraha (NPM: 714250007)"
    p.font.size = Pt(13)
    p.font.color.rgb = COLOR_WHITE
    p.space_before = Pt(4)

    p = tf_p.add_paragraph()
    p.text = "Program Studi D4 Teknik Informatika"
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_MUTED
    p.space_before = Pt(10)

    add_card(slide1, Inches(6.8), Inches(4.0), Inches(5.7), Inches(2.7), bg_color=RGBColor(35, 35, 60), border_color=COLOR_SECONDARY)
    tb_dos = slide1.shapes.add_textbox(Inches(7.0), Inches(4.1), Inches(5.3), Inches(2.5))
    tf_d = tb_dos.text_frame
    tf_d.word_wrap = True

    p = tf_d.paragraphs[0]
    p.text = "DOSEN PEMBIMBING & KOORDINATOR:"
    p.font.size = Pt(12)
    p.font.bold = True
    p.font.color.rgb = COLOR_SECONDARY

    p = tf_d.add_paragraph()
    p.text = "Pembimbing: Cahyo Prianto, S.Pd., M.T., CDSP, SFPC"
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_WHITE
    p.space_before = Pt(6)

    p = tf_d.add_paragraph()
    p.text = "Koordinator: M. Yusril Helmi Setyawan, S.Kom., M.Kom."
    p.font.size = Pt(12.5)
    p.font.color.rgb = COLOR_WHITE
    p.space_before = Pt(4)

    p = tf_d.add_paragraph()
    p.text = "Universitas Logistik dan Bisnis Internasional (ULBI) - 2026"
    p.font.size = Pt(11)
    p.font.color.rgb = COLOR_MUTED
    p.space_before = Pt(10)


    # ==================== SLIDE 2: LATAR BELAKANG ====================
    slide2 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(slide2, "1. LATAR BELAKANG PENGEMBANGAN")

    add_card(slide2, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2), title="KONTEKS AKADEMIK MAHASISWA ULBI", border_color=COLOR_PRIMARY)
    tb_lb = slide2.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.3), Inches(4.4))
    tf_lb = tb_lb.text_frame
    tf_lb.word_wrap = True

    items_lb = [
        ("Kesibukan Aktivitas Mahasiswa", "Mahasiswa D4 Teknik Informatika ULBI menghadapi padatnya aktivitas akademik mulai dari praktikum, tugas harian, dan tugas kelompok secara serentak."),
        ("Kelemahan Pencatatan Konvensional", "Metode pencatatan konvensional (buku fisik/notes hp biasa) bersifat pasif, tidak terpusat, dan tidak memiliki sistem pengingat otomatis yang proaktif."),
        ("Kebutuhan Monitoring Akademik", "Administrator akademik memerlukan sarana monitoring real-time untuk melihat perkembangan progres tugas dan kedisiplinan belajar mahasiswa."),
        ("Pemanfaatan Web & WA Gateway", "Laravel dan Tailwind CSS digunakan untuk membangun platform visual Kanban yang terintegrasi dengan WhatsApp Gateway Fonnte API sebagai asisten pengingat.")
    ]
    for i, (head, desc) in enumerate(items_lb):
        p = tf_lb.paragraphs[0] if i == 0 else tf_lb.add_paragraph()
        p.text = f"📌 {head}: "
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_PRIMARY
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(10.5)
        run.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(10) if i > 0 else Pt(0)


    # ==================== SLIDE 3: TUJUAN ====================
    slide3 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(slide3, "2. TUJUAN PENGEMBANGAN")

    add_card(slide3, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.2), title="TUJUAN UMUM", border_color=COLOR_PRIMARY)
    tb_tu = slide3.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.4))
    tf_tu = tb_tu.text_frame
    tf_tu.word_wrap = True

    p = tf_tu.paragraphs[0]
    p.text = "Membangun & mengimplementasikan aplikasi manajemen tugas berbasis web (TaskMate) menggunakan Laravel, Tailwind CSS, dan MySQL untuk meningkatkan efisiensi waktu, manajemen prioritas, serta produktivitas akademik mahasiswa ULBI."
    p.font.size = Pt(11.5)
    p.font.color.rgb = COLOR_TEXT

    add_card(slide3, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2), title="TUJUAN KHUSUS (TARGET LUARAN)", border_color=COLOR_SECONDARY)
    tb_tk = slide3.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.3), Inches(4.4))
    tf_tk = tb_tk.text_frame
    tf_tk.word_wrap = True

    tujuan_khusus = [
        "Otorisasi hak akses Multi-role (Admin vs Mahasiswa) terisolasi.",
        "Fitur Manajemen Tugas (CRUD) lengkap dengan prioritas & deadline.",
        "Modul Kanban Board visual 4 kolom (To Do, Doing, Review, Done).",
        "Integrasi WhatsApp Service API Fonnte untuk pengiriman naskah pengingat H-5 & H-2 (Manual Dispatch).",
        "Dasbor Monitoring Admin dengan statistik real-time & audit log."
    ]
    for i, tk in enumerate(tujuan_khusus):
        p = tf_tk.paragraphs[0] if i == 0 else tf_tk.add_paragraph()
        p.text = f"🎯 {tk}"
        p.font.size = Pt(10.5)
        p.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(8) if i > 0 else Pt(0)


    # ==================== SLIDE 4: MASALAH ====================
    slide4 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(slide4, "3. PERMASALAHAN UTAMA")

    add_card(slide4, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2), title="IDENTIFIKASI PERMASALAHAN AKADEMIK", border_color=COLOR_PRIMARY)
    tb_m = slide4.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(11.3), Inches(4.4))
    tf_m = tb_m.text_frame
    tf_m.word_wrap = True

    items_m = [
        ("Prokrastinasi & Kelupaan Batas Waktu", "Banyaknya tenggat waktu (due date) yang beragam memicu kelupaan tugas, penundaan pekerjaan, & penumpukan tugas di akhir semester."),
        ("Tidak Adanya Pengingat Otomatis", "Sistem pencatatan konvensional bersifat pasif & tidak dapat mengirimkan notifikasi pengingat otomatis ke handphone mahasiswa sebelum batas waktu habis."),
        ("Kesulitan Pengurutan Prioritas", "Mahasiswa kesulitan memvisualisasikan alur pengerjaan tugas akademik secara dinamis (tugas mana yang belum, sedang, atau selesai dikerjakan)."),
        ("Minimnya Sarana Pengawasan", "Pengelola akademik/Admin kesulitan memantau keaktifan tugas mahasiswa secara waktu nyata & tidak adanya log aktivitas sistem.")
    ]
    for i, (head, desc) in enumerate(items_m):
        p = tf_m.paragraphs[0] if i == 0 else tf_m.add_paragraph()
        p.text = f"❌ {head}: "
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_DARK
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(10.5)
        run.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(10) if i > 0 else Pt(0)


    # ==================== SLIDE 5: TAMPILAN APLIKASI WITH SCREENSHOTS ====================
    slide5 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(slide5, "4. TAMPILAN APLIKASI & MODUL UI/UX")

    # Left: descriptions card
    add_card(slide5, Inches(0.8), Inches(1.5), Inches(5.8), Inches(5.2), title="MODUL ANTARMUKA UTAMA", border_color=COLOR_PRIMARY)
    tb_ui_desc = slide5.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.4), Inches(4.4))
    tf_ui_desc = tb_ui_desc.text_frame
    tf_ui_desc.word_wrap = True

    ui_points = [
        ("Dashboard Mahasiswa", "Widget ringkasan data tugas aktif/selesai, progress bar kelulusan tugas, kalender tenggat, & feed notifikasi."),
        ("Interactive Kanban Board", "4 Kolom visual (To Do, Doing, Review, Done) dengan drag-and-drop SortableJS & auto-done 100% checklist."),
        ("Admin Control Panel", "Dasbor total mahasiswa, live board monitoring board mahasiswa, kelola akun (CRUD), dan log riwayat audit aktivitas.")
    ]
    for i, (m_title, m_desc) in enumerate(ui_points):
        p = tf_ui_desc.paragraphs[0] if i == 0 else tf_ui_desc.add_paragraph()
        p.text = f"📱 {m_title}: "
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_PRIMARY
        run = p.add_run()
        run.text = m_desc
        run.font.bold = False
        run.font.size = Pt(10.5)
        run.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(10) if i > 0 else Pt(0)

    # Right: screenshots stacked
    if os.path.exists(img_kanban):
        slide5.shapes.add_picture(img_kanban, Inches(6.9), Inches(1.5), width=Inches(5.6))
    if os.path.exists(img_admin):
        slide5.shapes.add_picture(img_admin, Inches(6.9), Inches(4.1), width=Inches(5.6))


    # ==================== SLIDE 6: DATABASE (SUNNAH) ====================
    slide6 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(slide6, "5. DATABASE MYSQL / SCHEMA ERD (SUNNAH)")

    add_card(slide6, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2))
    tb_db = slide6.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.3), Inches(4.8))
    tf_db = tb_db.text_frame
    tf_db.word_wrap = True

    tables_info = [
        ("mahasiswas (PK: id)", "nim (UK), nama, email (UK), password, phone, foto, universitas, semester"),
        ("tasks (PK: id, FK: mahasiswa_id)", "title, description, status (todo, doing, review, done), deadline, priority, checklist (JSON)"),
        ("activities (PK: id, FK: mahasiswa_id)", "activity_text, status_tugas, created_at (Audit trail log sistem)"),
        ("task_reminders (PK: id, FK: task_id, mahasiswa_id)", "days_before (5/2), sent_at (Catatan log status kirim WA Fonnte)"),
        ("admins (PK: id)", "nama, email (UK), password (Akun operator administrator)")
    ]

    p = tf_db.paragraphs[0]
    p.text = "STRUKTUR ENTITAS BASIS DATA & RELASI RELASIONAL:"
    p.font.size = Pt(13)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY

    for tname, fld in tables_info:
        p = tf_db.add_paragraph()
        p.text = f"📊 Tabel {tname}: "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_DARK
        run = p.add_run()
        run.text = fld
        run.font.bold = False
        run.font.size = Pt(10.5)
        run.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(8)


    # ==================== SLIDE 7: ALUR KERJA & BANTUAN HUBUNGI ADMIN (MANUAL WA REMINDER) ====================
    slide7 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(slide7, "6. ALUR KERJA APLIKASI & LAYANAN PUSAT BANTUAN")

    # Left Card: Operational workflow
    add_card(slide7, Inches(0.8), Inches(1.5), Inches(5.7), Inches(5.2), title="ALUR OPERASIONAL SYSTEM", border_color=COLOR_PRIMARY)
    tb_ak = slide7.shapes.add_textbox(Inches(1.0), Inches(2.1), Inches(5.3), Inches(4.4))
    tf_ak = tb_ak.text_frame
    tf_ak.word_wrap = True

    steps_list = [
        ("1. Registrasi & Login", "Mahasiswa daftar (Bcrypt) -> Login via Web Guard Auth ke workspace."),
        ("2. Input Tugas Baru", "Submit data tugas via AJAX -> Otomatis masuk status 'TO DO'."),
        ("3. Kanban & Auto-Done", "Drag & Drop tugas. Checklist 100% otomatis geser task ke status DONE."),
        ("4. Pengiriman Pengingat Manual", "Administrator memicu command PHP Artisan `tasks:send-deadline-reminders` untuk kirim pengingat tugas H-5 & H-2 deadline via Fonnte API.")
    ]
    for i, (head, desc) in enumerate(steps_list):
        p = tf_ak.paragraphs[0] if i == 0 else tf_ak.add_paragraph()
        p.text = f"🔄 {head}: "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_PRIMARY
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(9.5)
        run.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(6) if i > 0 else Pt(0)

    # Right Card: Help center (Bantuan Hubungi Admin)
    add_card(slide7, Inches(6.8), Inches(1.5), Inches(5.7), Inches(5.2), title="BANTUAN UNTUK MENGHUBUNGI ADMIN", border_color=COLOR_SECONDARY)
    tb_help = slide7.shapes.add_textbox(Inches(7.0), Inches(2.1), Inches(5.3), Inches(4.4))
    tf_help = tb_help.text_frame
    tf_help.word_wrap = True

    help_points = [
        ("Modul Pusat Bantuan Mahasiswa", "Mahasiswa dapat mengakses halaman Bantuan (Help Center) langsung dari dashboard mereka jika menemui kendala teknis."),
        ("WhatsApp Hubungi Admin", "Sistem menyediakan tombol khusus terintegrasi WhatsApp (Direct API). Mahasiswa dapat mengirimkan pesan aduan kendala dengan teks template otomatis."),
        ("Respon Cepat & Solutif", "Mempermudah komunikasi dua arah antara Mahasiswa dan Admin ULBI untuk verifikasi bug, reset sandi, atau panduan penggunaan sistem.")
    ]
    for i, (h_title, h_desc) in enumerate(help_points):
        p = tf_help.paragraphs[0] if i == 0 else tf_help.add_paragraph()
        p.text = f"💬 {h_title}: "
        p.font.bold = True
        p.font.size = Pt(11)
        p.font.color.rgb = COLOR_SECONDARY
        run = p.add_run()
        run.text = h_desc
        run.font.bold = False
        run.font.size = Pt(9.5)
        run.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(10) if i > 0 else Pt(0)


    # ==================== SLIDE 8: HASIL & SIMPULAN ====================
    slide8 = prs.slides.add_slide(blank_layout)
    add_slide_header_and_footer(slide8, "7. HASIL AKHIR & KESIMPULAN")

    add_card(slide8, Inches(0.8), Inches(1.5), Inches(11.7), Inches(5.2), border_color=COLOR_PRIMARY)
    tb_res = slide8.shapes.add_textbox(Inches(1.0), Inches(1.7), Inches(11.3), Inches(4.8))
    tf_res = tb_res.text_frame
    tf_res.word_wrap = True

    kesimpulan = [
        ("Implementasi Fungsional 100%", "Seluruh target luaran sistem (Multi-role Auth, CRUD tugas, Kanban Board, Scheduler, WA Gateway) telah sukses dibangun."),
        ("Peningkatan Produktivitas & Disiplin", "Papan Kanban memvisualisasikan alur tugas dengan jelas, sementara WhatsApp Gateway proaktif menekan angka kelupaan/keterlambatan tugas."),
        ("Transparansi Pengawasan", "Dasbor admin memberikan kemudahan bagi pengelola akademik untuk memantau progres tugas mahasiswa serta mencatat riwayat log audit aktivitas."),
        ("Rencana Pengembangan", "Ke depannya dapat ditingkatkan dengan bot WhatsApp interaktif dua arah, grafik analitik produktivitas mingguan, & Progressive Web App (PWA).")
    ]

    for i, (head, desc) in enumerate(kesimpulan):
        p = tf_res.paragraphs[0] if i == 0 else tf_res.add_paragraph()
        p.text = f"✅ {i+1}. {head}: "
        p.font.bold = True
        p.font.size = Pt(11.5)
        p.font.color.rgb = COLOR_PRIMARY
        run = p.add_run()
        run.text = desc
        run.font.bold = False
        run.font.size = Pt(10.5)
        run.font.color.rgb = COLOR_TEXT
        p.space_before = Pt(12) if i > 0 else Pt(0)


    # ==================== SLIDE 9: Q&A / CLOSING ====================
    slide9 = prs.slides.add_slide(blank_layout)
    bg9 = slide9.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg9.fill.solid()
    bg9.fill.fore_color.rgb = COLOR_DARK
    bg9.line.fill.background()

    tb_q = slide9.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(11.3), Inches(3.5))
    tf_q = tb_q.text_frame
    tf_q.word_wrap = True

    p = tf_q.paragraphs[0]
    p.text = "TERIMA KASIH"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLOR_PRIMARY
    p.alignment = PP_ALIGN.CENTER

    p2 = tf_q.add_paragraph()
    p2.text = "Sesi Tanya Jawab (Q & A) & Demo Live Aplikasi"
    p2.font.size = Pt(24)
    p2.font.color.rgb = COLOR_WHITE
    p2.alignment = PP_ALIGN.CENTER
    p2.space_before = Pt(10)

    p3 = tf_q.add_paragraph()
    p3.text = "TaskMate: Teman Belajar & Solusi Manajemen Tugas Mahasiswa\nUniversitas Logistik dan Bisnis Internasional (ULBI)"
    p3.font.size = Pt(14)
    p3.font.color.rgb = COLOR_SECONDARY
    p3.alignment = PP_ALIGN.CENTER
    p3.space_before = Pt(20)

    output_path = "TaskMate_Presentasi_Sidang_V3.pptx"
    prs.save(output_path)
    print(f"Presentation updated with manual WA reminder details and saved to {os.path.abspath(output_path)}")

if __name__ == "__main__":
    create_presentation()
